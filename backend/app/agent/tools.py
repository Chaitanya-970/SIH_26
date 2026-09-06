import asyncio
import os
import uuid
from collections.abc import Callable, Awaitable
from dataclasses import dataclass
from typing import Any
import logging
from anyio import to_thread

import chromadb # type: ignore
import pandas as pd # type: ignore
from pypdf import PdfReader # type: ignore
from docx import Document as DocxDocument # type: ignore
from pptx import Presentation # type: ignore

from app.models.schemas import ToolCall, ToolResult
from app.config import Settings, ModelRegistry
from app.services.ollama import OllamaClient

logger = logging.getLogger(__name__)

@dataclass
class ToolContext:
    session_id: str
    uploads_dir: str
    exports_dir: str
    ollama: OllamaClient
    registry: ModelRegistry
    settings: Settings

async def execute_code(args: dict[str, Any], context: ToolContext) -> ToolResult:
    script = args.get("script", "")
    if not script.strip():
        return ToolResult(tool="execute_code", result="Error: empty script", success=False)

    timeout = context.settings.sandbox_timeout_seconds

    try:
        # Primary: Docker sandbox (when Vedant's image is ready)
        # Fallback: direct subprocess (for development/testing)
        proc = await asyncio.create_subprocess_exec(
            "docker", "run", "--rm",
            "--network", "none",
            "--memory", "512m",
            "--cpus", "1",
            "-i",  # read script from stdin
            "citadel-sandbox",
            "python", "-c", script,
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
        )
        stdout, stderr = await asyncio.wait_for(
            proc.communicate(), timeout=timeout
        )
    except FileNotFoundError:
        # Docker not available — fallback to direct python
        logger.warning("Docker not available — running code in UNSANDBOXED fallback mode")
        proc = await asyncio.create_subprocess_exec(
            "python", "-c", script,
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
        )
        stdout, stderr = await asyncio.wait_for(
            proc.communicate(), timeout=timeout
        )
    except asyncio.TimeoutError:
        try:
            proc.kill()
        except Exception:
            pass
        return ToolResult(
            tool="execute_code",
            result=f"Execution timed out after {timeout} seconds.",
            success=False
        )

    output = stdout.decode("utf-8", errors="replace") if stdout else ""
    errors = stderr.decode("utf-8", errors="replace") if stderr else ""

    result_text = ""
    if output:
        result_text += f"stdout:\n{output}"
    if errors:
        result_text += f"\nstderr:\n{errors}"
    if not result_text:
        result_text = "(no output)"

    # if process exited cleanly but is just printing something, it's a success
    # Note: proc.returncode is an Optional[int], but after communicate() it is an int.
    # We cast to avoid mypy issues if needed, but proc.returncode == 0 is safe.
    success = proc.returncode == 0

    return ToolResult(
        tool="execute_code",
        result=result_text.strip(),
        success=success
    )

async def search_knowledge_base(args: dict[str, Any], context: ToolContext) -> ToolResult:
    query = args.get("query", "")
    if not query.strip():
        return ToolResult(tool="search_knowledge_base", result="Error: empty query", success=False)

    # Get embedding model from registry
    embed_model = context.registry.get_by_capability("embedding")
    if not embed_model:
        return ToolResult(tool="search_knowledge_base", result="Error: no embedding model configured", success=False)

    try:
        # Embed the query
        query_embedding = await context.ollama.embed(embed_model.ollama_tag, query)

        # Query ChromaDB
        client = chromadb.PersistentClient(path=context.settings.chroma_dir)
        collection = client.get_or_create_collection(
            name="knowledge_base",
            metadata={"hnsw:space": "cosine"}
        )

        if collection.count() == 0:
            return ToolResult(
                tool="search_knowledge_base",
                result="Knowledge base is empty. No documents have been uploaded yet.",
                success=True
            )

        results = collection.query(
            query_embeddings=[query_embedding], # type: ignore[arg-type]
            n_results=3,
            include=["documents", "metadatas", "distances"]
        )

        # Format results
        chunks = []
        docs = results.get("documents")
        metas = results.get("metadatas")
        dists = results.get("distances")
        if docs and metas and dists:
            for doc, meta, dist in zip(docs[0], metas[0], dists[0]):
                # Handle possible None types in ChromaDB result lists
                doc_text = doc if doc else ""
                meta_dict = meta if meta else {}
                dist_val = dist if dist is not None else 1.0
                
                source = meta_dict.get("source", "unknown")
                relevance = 1.0 - dist_val
                chunks.append(f"[Source: {source} | Relevance: {relevance:.2f}]\n{doc_text}")

        return ToolResult(
            tool="search_knowledge_base",
            result="\n\n---\n\n".join(chunks) if chunks else "No relevant documents found.",
            success=True
        )
    except Exception as e:
        logger.warning(f"Error in search_knowledge_base: {e}")
        return ToolResult(tool="search_knowledge_base", result=f"Search error: {str(e)}", success=False)

def _sync_read_document(full_path: str, ext: str) -> str:
    if ext == ".pdf":
        reader = PdfReader(full_path)
        return "\n\n".join(page.extract_text() or "" for page in reader.pages)
    elif ext == ".docx":
        doc = DocxDocument(full_path)
        return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())
    elif ext in (".xlsx", ".xls"):
        df = pd.read_excel(full_path)
        return df.to_string(index=False)
    elif ext == ".csv":
        df = pd.read_csv(full_path)
        return df.to_string(index=False)
    elif ext == ".pptx":
        prs = Presentation(full_path)
        all_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    all_text.append(shape.text)
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                all_text.append(slide.notes_slide.notes_text_frame.text)
        return "\n\n".join(all_text).strip()
    else:
        with open(full_path, "r", encoding="utf-8", errors="replace") as f:
            return f.read()

async def read_document(args: dict[str, Any], context: ToolContext) -> ToolResult:
    file_path = args.get("file_path", "")
    if not file_path:
        return ToolResult(tool="read_document", result="Error: no file_path provided", success=False)

    # Resolve to session uploads directory (prevent path traversal)
    base_name = os.path.basename(file_path)
    full_path = os.path.normpath(os.path.join(context.uploads_dir, base_name))
    
    # Ensure it's still inside uploads_dir (note: os.path.basename already neutralizes traversal, this is defense-in-depth against symlinks)
    if not full_path.startswith(os.path.normpath(context.uploads_dir)):
        return ToolResult(tool="read_document", result="Error: invalid file path", success=False)

    if not os.path.exists(full_path):
        return ToolResult(tool="read_document", result=f"Error: file not found: {file_path}", success=False)

    ext = os.path.splitext(full_path)[1].lower()

    try:
        text = await to_thread.run_sync(_sync_read_document, full_path, ext)

        # Truncate very long documents to fit within context
        max_chars = 8000
        if len(text) > max_chars:
            text = text[:max_chars] + f"\n\n[... truncated, showing first {max_chars} chars of {len(text)} total]"

        return ToolResult(tool="read_document", result=text, success=True)
    except Exception as e:
        logger.warning(f"Error reading document {file_path}: {e}")
        return ToolResult(tool="read_document", result=f"Read error: {str(e)}", success=False)

def _sync_write_word_document(title: str, body: str, sig_block: str, exports_dir: str) -> str:
    filename = title.lower().replace(" ", "_")[:50] + ".docx"
    filepath = os.path.join(exports_dir, filename)
    doc = DocxDocument()
    doc.add_heading(title, level=1)
    for paragraph in body.split("\n"):
        if paragraph.strip():
            doc.add_paragraph(paragraph.strip())
    if sig_block:
        doc.add_paragraph("")  # spacer
        doc.add_paragraph(sig_block)
    os.makedirs(exports_dir, exist_ok=True)
    doc.save(filepath)
    return filename

async def write_word_document(args: dict[str, Any], context: ToolContext) -> ToolResult:
    content = args.get("content")
    if not isinstance(content, dict):
        return ToolResult(tool="write_word_document", result="Error: 'content' must be a dict with title/body keys", success=False)
        
    title = content.get("title", "Untitled Document")
    body = content.get("body", "")
    sig_block = content.get("signature_block", "")

    try:
        filename = await to_thread.run_sync(_sync_write_word_document, title, body, sig_block, context.exports_dir)
        return ToolResult(
            tool="write_word_document",
            result=f"Document saved: {filename}",
            success=True,
            filename=filename
        )
    except Exception as e:
        logger.warning(f"Error writing word document: {e}")
        return ToolResult(tool="write_word_document", result=f"Error generating document: {str(e)}", success=False)

def _sync_write_spreadsheet(data: list[Any], exports_dir: str) -> tuple[str, int, int]:
    filename = f"export_{uuid.uuid4().hex[:8]}.xlsx"
    filepath = os.path.join(exports_dir, filename)
    df = pd.DataFrame(data)
    os.makedirs(exports_dir, exist_ok=True)
    df.to_excel(filepath, index=False, engine="openpyxl")
    return filename, len(df), len(df.columns)

async def write_spreadsheet(args: dict[str, Any], context: ToolContext) -> ToolResult:
    data = args.get("data", [])
    if not data or not isinstance(data, list):
        return ToolResult(tool="write_spreadsheet", result="Error: 'data' must be a non-empty list of objects", success=False)

    try:
        filename, rows, cols = await to_thread.run_sync(_sync_write_spreadsheet, data, context.exports_dir)
        return ToolResult(
            tool="write_spreadsheet",
            result=f"Spreadsheet saved: {filename} ({rows} rows, {cols} columns)",
            success=True
        )
    except Exception as e:
        logger.warning(f"Error writing spreadsheet: {e}")
        return ToolResult(tool="write_spreadsheet", result=f"Error generating spreadsheet: {str(e)}", success=False)

def _sync_write_presentation(title: str, slides_data: list[Any], exports_dir: str) -> str:
    filename = title.lower().replace(" ", "_")[:50] + ".pptx"
    filepath = os.path.join(exports_dir, filename)
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    if slide.shapes.title:
        slide.shapes.title.text = title
    if len(slide.placeholders) > 1 and slide.placeholders[1]:
        slide.placeholders[1].text = "Generated by CITADEL WORKSPACE"
    for slide_data in slides_data:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        if slide.shapes.title:
            slide.shapes.title.text = slide_data.get("heading", "")
        if len(slide.placeholders) > 1 and slide.placeholders[1]:
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.clear()
            for i, bullet in enumerate(slide_data.get("bullets", [])):
                if i == 0:
                    tf.text = bullet
                else:
                    p = tf.add_paragraph()
                    p.text = bullet
    os.makedirs(exports_dir, exist_ok=True)
    prs.save(filepath)
    return filename

async def write_presentation(args: dict[str, Any], context: ToolContext) -> ToolResult:
    content = args.get("content")
    if not isinstance(content, dict):
        return ToolResult(tool="write_presentation", result="Error: 'content' must be a dict with title/slides keys", success=False)
        
    title = content.get("title", "Untitled Presentation")
    slides_data = content.get("slides", [])

    try:
        filename = await to_thread.run_sync(_sync_write_presentation, title, slides_data, context.exports_dir)
        return ToolResult(
            tool="write_presentation",
            result=f"Presentation saved: {filename} ({len(slides_data)} content slides)",
            success=True
        )
    except Exception as e:
        logger.warning(f"Error writing presentation: {e}")
        return ToolResult(tool="write_presentation", result=f"Error generating presentation: {str(e)}", success=False)

TOOL_REGISTRY: dict[str, Callable[[dict[str, Any], ToolContext], Awaitable[ToolResult]]] = {
    "execute_code": execute_code,
    "search_knowledge_base": search_knowledge_base,
    "read_document": read_document,
    "write_word_document": write_word_document,
    "write_spreadsheet": write_spreadsheet,
    "write_presentation": write_presentation,
}

async def dispatch_tool(tool_call: ToolCall, context: ToolContext) -> ToolResult:
    func = TOOL_REGISTRY.get(tool_call.tool)
    if not func:
        logger.warning(f"Unknown tool called: {tool_call.tool}")
        return ToolResult(tool=tool_call.tool, result=f"Unknown tool: {tool_call.tool}", success=False)
    try:
        return await func(tool_call.args, context)
    except Exception as e:
        logger.error(f"Unhandled error in tool {tool_call.tool}: {e}")
        return ToolResult(tool=tool_call.tool, result=f"Tool error: {str(e)}", success=False)
